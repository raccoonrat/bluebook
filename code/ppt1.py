import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_LINE_DASH_STYLE

# ==================== 颜色常量定义 ====================
COLOR_TRACK_A_HEADER = RGBColor(78, 151, 204)
COLOR_TRACK_A_BG = RGBColor(221, 238, 248)
COLOR_TRACK_B_HEADER = RGBColor(243, 154, 122)
COLOR_TRACK_B_BG = RGBColor(249, 227, 216)
COLOR_SHARED_BACKBONE_BG = RGBColor(232, 232, 232)
COLOR_SHARED_BACKBONE_BORDER = RGBColor(110, 110, 110)
COLOR_GOVERNANCE_BG_OUTER = RGBColor(185, 227, 199)
COLOR_GOVERNANCE_BG_INNER = RGBColor(223, 242, 228)
COLOR_SCOPE_CARD = RGBColor(191, 216, 245)
COLOR_EVIDENCE_CARD = RGBColor(201, 216, 232)
COLOR_CONTRACT_CARD = RGBColor(248, 211, 178)
COLOR_DECISION_CARD = RGBColor(243, 168, 143)
COLOR_RELEASE_CARD = RGBColor(157, 206, 154)
COLOR_REDUCED_CAP_CARD = RGBColor(246, 219, 124)
COLOR_CANARY_CARD = RGBColor(148, 199, 242)
COLOR_ROLLBACK_CARD = RGBColor(233, 125, 125)
COLOR_TEXT_MAIN = RGBColor(17, 17, 17)
COLOR_BORDER_DARK = RGBColor(51, 51, 51)
COLOR_ARROW_GRAY = RGBColor(207, 207, 207)
COLOR_ARROW_BLUE = RGBColor(91, 155, 213)
COLOR_ARROW_ORANGE = RGBColor(244, 162, 97)
COLOR_DASHED_LINE = RGBColor(100, 100, 100)

# ==================== 工具函数 ====================
def add_round_rect(slide, left, top, width, height, fill_color, line_color, line_width=Pt(1.8)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = line_width
    shape.line.dash_style = MSO_LINE_DASH_STYLE.SOLID
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=Pt(10), font_bold=False, align=PP_ALIGN.CENTER, text_color=COLOR_TEXT_MAIN):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = font_bold
    p.alignment = align
    p.font.color.rgb = text_color
    return txBox

def add_block_arrow(slide, left, top, width, height, fill_color, line_color):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = fill_color
    arrow.line.color.rgb = line_color
    arrow.line.width = Pt(1.5)
    return arrow

def add_line(slide, x1, y1, x2, y2, color=COLOR_BORDER_DARK, width=Pt(1), dashed=False):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = width
    if dashed:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line

def add_icon_document(slide, left, top, size=Inches(0.6)):
    doc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, size, size*0.8)
    doc.fill.solid()
    doc.fill.fore_color.rgb = RGBColor(255,255,255)
    doc.line.color.rgb = COLOR_BORDER_DARK
    for i in range(3):
        y = top + Pt(8) + i*Pt(12)
        line = add_line(slide, left+Pt(4), y, left+size-Pt(4), y)
    return doc

def add_icon_shield(slide, left, top, size=Inches(0.7)):
    shield = add_round_rect(slide, left, top, size, size, RGBColor(255,255,255), COLOR_BORDER_DARK, Pt(1))
    cx, cy = left+size/2, top+size/2
    v_line = add_line(slide, cx, top+Pt(4), cx, top+size-Pt(4))
    h_line = add_line(slide, left+Pt(4), cy, left+size-Pt(4), cy)
    add_textbox(slide, left+Pt(4), top+Pt(4), size/2-Pt(8), size/2-Pt(8), "Security", Pt(7), True)
    add_textbox(slide, cx+Pt(4), top+Pt(4), size/2-Pt(8), size/2-Pt(8), "Utility", Pt(7), True)
    add_textbox(slide, left+Pt(4), cy+Pt(4), size/2-Pt(8), size/2-Pt(8), "System", Pt(7), True)
    add_textbox(slide, cx+Pt(4), cy+Pt(4), size/2-Pt(8), size/2-Pt(8), "Governance", Pt(7), True)
    return shield

def add_icon_folder(slide, left, top, size=Inches(0.6)):
    tab = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, size*0.5, Inches(0.12))
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top+Inches(0.12), size, size*0.7)
    for s in [tab, body]:
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(245,245,245)
        s.line.color.rgb = COLOR_BORDER_DARK
    return body

def add_icon_stamp(slide, left, top, size=Inches(0.6)):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top+size*0.2, size, size*0.6)
    handle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left+size*0.4, top, size*0.2, size*0.2)
    for s in [circle, handle]:
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(255,255,255)
        s.line.color.rgb = COLOR_BORDER_DARK
    return circle

def add_icon_server(slide, left, top, size=Inches(0.7)):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, size, size)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(240,240,240)
    box.line.color.rgb = COLOR_BORDER_DARK
    for i in range(3):
        y = top + size*0.2 + i*size*0.25
        slot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left+Pt(4), y, size-Pt(8), size*0.15)
        slot.line.color.rgb = COLOR_BORDER_DARK
    return box

def add_icon_database(slide, left, top, size=Inches(0.7)):
    top_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size*0.3)
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top+size*0.15, size, size*0.6)
    bottom_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top+size*0.6, size, size*0.3)
    for s in [top_oval, body, bottom_oval]:
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(245,245,245)
        s.line.color.rgb = COLOR_BORDER_DARK
    return body

def add_icon_bird(slide, left, top, size=Inches(0.4)):
    body = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size*0.6)
    wing = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left+size*0.2, top+size*0.2, size*0.3, size*0.3)
    body.fill.solid()
    wing.fill.solid()
    body.fill.fore_color.rgb = wing.fill.fore_color.rgb = COLOR_CANARY_CARD
    body.line.color.rgb = wing.line.color.rgb = COLOR_BORDER_DARK
    return body

def add_icon_rollback(slide, left, top, size=Inches(0.5)):
    arrow = slide.shapes.add_shape(MSO_SHAPE.CIRCULAR_ARROW, left, top, size, size)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_ROLLBACK_CARD
    arrow.line.color.rgb = COLOR_BORDER_DARK
    return arrow

# ==================== 主绘图函数 ====================
def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = Inches(13.333), Inches(7.281)

    # Track A
    track_a = add_round_rect(slide, Inches(0.5), Inches(0.5), Inches(4.2), Inches(3.0), COLOR_TRACK_A_BG, COLOR_BORDER_DARK)
    add_round_rect(slide, Inches(0.5), Inches(0.5), Inches(4.2), Inches(0.4), COLOR_TRACK_A_HEADER, COLOR_BORDER_DARK)
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(4.2), Inches(0.4), "Track A", Pt(14), True, PP_ALIGN.CENTER, RGBColor(255,255,255))
    add_icon_document(slide, Inches(0.7), Inches(1.0))
    add_textbox(slide, Inches(0.6), Inches(1.6), Inches(0.9), Inches(0.3), "Deployable\nConclusions", Pt(9), True)
    add_icon_shield(slide, Inches(1.7), Inches(1.0))
    add_textbox(slide, Inches(1.6), Inches(1.6), Inches(0.9), Inches(0.3), "Metrics board", Pt(9), True)
    add_icon_folder(slide, Inches(2.7), Inches(1.0))
    add_textbox(slide, Inches(2.6), Inches(1.6), Inches(0.9), Inches(0.3), "Acceptance\nContract", Pt(9), True)

    # Track A Arrow
    add_block_arrow(slide, Inches(4.8), Inches(1.8), Inches(1.0), Inches(0.5), COLOR_ARROW_BLUE, COLOR_BORDER_DARK)
    add_textbox(slide, Inches(4.8), Inches(1.6), Inches(1.0), Inches(0.2), "Track A Evidence Flow", Pt(8), True)

    # Track B
    track_b = add_round_rect(slide, Inches(0.5), Inches(3.8), Inches(4.2), Inches(3.0), COLOR_TRACK_B_BG, COLOR_BORDER_DARK)
    add_round_rect(slide, Inches(0.5), Inches(3.8), Inches(4.2), Inches(0.4), COLOR_TRACK_B_HEADER, COLOR_BORDER_DARK)
    add_textbox(slide, Inches(0.5), Inches(3.8), Inches(4.2), Inches(0.4), "Track B", Pt(14), True, PP_ALIGN.CENTER, RGBColor(255,255,255))
    add_textbox(slide, Inches(1.5), Inches(4.5), Inches(2.0), Inches(0.3), "No Mixing Into Track A", Pt(10), True, text_color=RGBColor(200,0,0))

    # Track B Arrow
    add_block_arrow(slide, Inches(4.8), Inches(5.0), Inches(1.0), Inches(0.5), COLOR_ARROW_ORANGE, COLOR_BORDER_DARK)
    add_textbox(slide, Inches(4.8), Inches(4.8), Inches(1.0), Inches(0.2), "Track B Mechanism Flow", Pt(8), True)

    # Shared Backbone
    backbone = add_round_rect(slide, Inches(6.0), Inches(0.5), Inches(1.8), Inches(6.3), COLOR_SHARED_BACKBONE_BG, COLOR_SHARED_BACKBONE_BORDER)
    add_textbox(slide, Inches(6.0), Inches(0.5), Inches(1.8), Inches(0.4), "Shared Backbone", Pt(11), True)
    add_textbox(slide, Inches(6.0), Inches(1.0), Inches(1.8), Inches(0.3), "CI/CD pipeline", Pt(9), True)
    add_icon_server(slide, Inches(6.3), Inches(1.4))
    add_textbox(slide, Inches(6.0), Inches(2.3), Inches(1.8), Inches(0.3), "Runtime Execution", Pt(9), True)
    add_icon_database(slide, Inches(6.3), Inches(4.8))

    # Governance Board
    gov = add_round_rect(slide, Inches(8.2), Inches(0.5), Inches(3.8), Inches(5.0), COLOR_GOVERNANCE_BG_OUTER, COLOR_BORDER_DARK)
    add_textbox(slide, Inches(8.2), Inches(0.5), Inches(3.8), Inches(0.4), "Unified Governance Board", Pt(12), True)
    add_round_rect(slide, Inches(8.5), Inches(1.0), Inches(3.2), Inches(4.0), COLOR_GOVERNANCE_BG_INNER, COLOR_BORDER_DARK)
    add_textbox(slide, Inches(8.7), Inches(1.1), Inches(2.8), Inches(0.3), "Evidence-Referenced Decision", Pt(10), True)

    # Decision Cards
    decision = add_round_rect(slide, Inches(12.3), Inches(0.5), Inches(1.5), Inches(5.5), RGBColor(245,245,245), COLOR_BORDER_DARK)
    add_textbox(slide, Inches(12.3), Inches(0.5), Inches(1.5), Inches(0.4), "Decision Cards", Pt(11), True)
    add_round_rect(slide, Inches(12.5), Inches(1.0), Inches(1.1), Inches(0.9), COLOR_RELEASE_CARD, COLOR_BORDER_DARK)
    add_textbox(slide, Inches(12.5), Inches(1.2), Inches(1.1), Inches(0.5), "Release card", Pt(9), True)
    add_round_rect(slide, Inches(12.5), Inches(2.1), Inches(1.1), Inches(0.9), COLOR_REDUCED_CAP_CARD, COLOR_BORDER_DARK)
    add_textbox(slide, Inches(12.5), Inches(2.3), Inches(1.1), Inches(0.5), "Reduced-capability", Pt(8), True)
    add_round_rect(slide, Inches(12.5), Inches(3.2), Inches(1.1), Inches(0.9), COLOR_CANARY_CARD, COLOR_BORDER_DARK)
    add_icon_bird(slide, Inches(12.8), Inches(3.4))
    add_textbox(slide, Inches(12.5), Inches(3.8), Inches(1.1), Inches(0.3), "Canary", Pt(9), True)
    add_round_rect(slide, Inches(12.5), Inches(4.3), Inches(1.1), Inches(0.9), COLOR_ROLLBACK_CARD, COLOR_BORDER_DARK)
    add_icon_rollback(slide, Inches(12.8), Inches(4.4))
    add_textbox(slide, Inches(12.5), Inches(4.9), Inches(1.1), Inches(0.3), "Rollback", Pt(9), True)

    # Feedback Lines
    add_line(slide, Inches(2.5), Inches(6.8), Inches(6.5), Inches(6.8), COLOR_DASHED_LINE, Pt(1), dashed=True)
    add_textbox(slide, Inches(4.0), Inches(6.6), Inches(1.5), Inches(0.2), "Mechanism Feedback", Pt(8), True)

# ==================== 主函数 ====================
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.281)
    build_slide(prs)
    prs.save("recreate_figure.pptx")
    print("✅ 生成完成：recreate_figure.pptx")

if __name__ == "__main__":
    main()